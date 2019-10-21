from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

class Present():

    def test_pptx(self):
        # create presentation with 1 slide ------
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = ['East', 'West', 'Midwest']
        chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

        prs.save('../presentation-tests/chart-01.pptx')

    def test_pptx02(self):
        # create presentation with 1 slide ------
        prs = Presentation("../presentation-tests/AIA Dallas 5Apr18_STD209.pptx")
        text_runs = []

        for indx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if indx == 21:
                            text_runs.append("{}-{}".format(indx, run.text))
                            run.text = run.text + "-modified"
        print(text_runs)
        prs.save('../presentation-tests/AIA Dallas 5Apr18_STD209-modified.pptx')

    def test_pptx03(self):
        # create presentation with 1 slide ------
        prs = Presentation("../presentation-tests/AIA Dallas 5Apr18_STD209.pptx")
        text_runs = []

        for slide in prs.slides:
            if not slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = 'new notes'
            else:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = 'next text on existing notes slide'
        for slide in prs.slides:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            p = text_frame.add_paragraph()
            p.text = 'next line added'
        prs.save('../presentation-tests/AIA Dallas 5Apr18_STD209-addedNotes.pptx')
